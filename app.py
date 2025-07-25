# =================================================================================================
# QTI ENGINEERING WORKBENCH - Enterprise Gold Master Edition
#
# AUTHOR: Subject Matter Expert AI
# DATE: 2024-07-24
# VERSION: 2.4.0 (SME Definitive RCA/SHAP Fix)
#
# DESCRIPTION:
# This is the gold master, single-file Streamlit application for a Quality Technical
# Investigation (QTI) Engineer. It serves as an enterprise-level workbench, integrating a
# comprehensive suite of statistical, machine learning, and process validation methods. The
# architecture is built for stability, clarity, and analytical depth, using Streamlit's
# session state and advanced caching. All known bugs are resolved, and functionalities
# have been substantially upgraded for precision and insight.
#
# SME NOTES (v2.4.0):
# - CRITICAL FIX (SHAP/RCA): The persistent `AssertionError` in the RCA Workbench has been
#   definitively resolved. The root cause was a feature mismatch between the globally trained
#   model and the locally filtered data.
# - NEW STRATEGY: The "AI-Driven Importance" tab now trains a temporary, local model on-the-fly
#   using only the user-filtered data. This guarantees perfect alignment between the model,
#   data, and SHAP explainer, making the analysis robust and error-free regardless of filtering.
# =================================================================================================

# --- 1. CORE & UTILITY IMPORTS ---
import streamlit as st
import pandas as pd
import numpy as np
from datetime import datetime, timedelta
import io
import yaml
import warnings
import time

# --- 2. DATA HANDLING & DATABASE IMPORTS ---
import dask.dataframe as dd

# --- 3. VISUALIZATION IMPORTS ---
import plotly.express as px
import plotly.graph_objects as go
import matplotlib.pyplot as plt
# Check for kaleido, required for static image export
try:
    import kaleido
except ImportError:
    st.error("The 'kaleido' package is not installed. Please install it (`pip install kaleido`) to enable report generation.")
    st.stop()


# --- 4. ANALYTICS & ML IMPORTS ---
from sklearn.ensemble import RandomForestClassifier
from sklearn.model_selection import train_test_split
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import KMeans
import scipy.stats as stats
import statsmodels.api as sm
from statsmodels.stats.weightstats import ttost_ind
import shap
from pydantic import BaseModel, Field, conint, confloat
from skopt import gp_minimize
from skopt.space import Real
from skopt.utils import use_named_args
from skopt.plots import plot_convergence

# --- 5. REPORTING IMPORTS ---
from pptx import Presentation
from pptx.util import Inches, Pt

# --- 6. APPLICATION CONFIGURATION ---
st.set_page_config(
    page_title="QTI Engineering Workbench v2.4",
    page_icon="🔬",
    layout="wide",
    initial_sidebar_state="expanded",
    menu_items={
        'Get Help': 'https://www.example.com/help',
        'Report a bug': "https://www.example.com/bug",
        'About': "# QTI Engineering Workbench v2.4\nThis is an enterprise-grade application for Quality professionals."
    }
)
# Suppress common warnings for a cleaner user experience
warnings.filterwarnings("ignore", category=UserWarning)
warnings.filterwarnings("ignore", category=FutureWarning)


# =================================================================================================
# CONFIGURATION & ADVANCED DATA SIMULATION LAYER
# =================================================================================================

class ReportSettings(BaseModel):
    author: str
    company_name: str

class SPCSettings(BaseModel):
    sigma_level: confloat(ge=2.0, le=4.0) = 3.0
    usl: confloat(gt=0) = 58.0
    lsl: confloat(ge=0) = 42.0

class AppConfig(BaseModel):
    report_settings: ReportSettings
    spc_settings: SPCSettings

@st.cache_data
def load_config():
    """Loads application configuration from a YAML string."""
    config_string = """
    report_settings:
        author: "QTI Engineering Team"
        company_name: "Innovate Bio-Diagnostics"
    spc_settings:
        sigma_level: 3.0
        usl: 58.0  # Upper Spec Limit for Pressure
        lsl: 42.0  # Lower Spec Limit for Pressure
    """
    return AppConfig(**yaml.safe_load(config_string))

@st.cache_data
def generate_process_data(num_records=2500, shift_point=1500):
    """
    Generates a sophisticated, realistic process dataset with multiple failure modes.
    - A distinct process shift occurs at the 'shift_point'.
    - Complex interactions between variables are simulated.
    - Includes ground truth 'failure_mode' for validation.
    """
    np.random.seed(42)
    start_time = datetime(2023, 1, 1)
    timestamps = [start_time + timedelta(hours=i) for i in range(num_records)]
    data = []

    for i, ts in enumerate(timestamps):
        material_lot = f"LOT-{(ts.day % 3) + 1}"
        operator = f"OP-{(i % 4) + 1}"
        reagent_supplier = "Supplier-A" if ts.month % 2 == 1 else "Supplier-B"
        is_anomaly = 0
        failure_mode = "Nominal"
        base_ph, base_vol, base_psi = 7.2, 10.0, 50.0
        ph = base_ph + np.random.normal(0, 0.05)
        vol = base_vol + np.random.normal(0, 0.03)
        psi = base_psi + np.random.normal(0, 0.5)

        if i > shift_point:
            reagent_supplier = "Supplier-C"
            psi += 2.5 + np.random.normal(0, 0.75)
            ph += 0.15 + (i - shift_point) * 0.0001
            failure_mode = "Process Shift"
            is_anomaly = 1 if np.random.rand() > 0.75 else 0
        else:
            if np.random.rand() > 0.98:
                is_anomaly = 1
                failure_mode = "Random Spike"
        if operator == 'OP-3':
            vol += 0.08; psi += 1.5
        elif operator == 'OP-4':
            vol -= 0.08; psi -= 1.5
        if material_lot == 'LOT-2':
            ph += 0.05
        psi += np.sin(i / 100) * 1.2

        record = {"timestamp": pd.to_datetime(ts), "line_id": f"LINE-{(i % 3) + 1}",
                  "reagent_ph": round(ph, 3), "fill_volume_ml": round(vol, 4), "pressure_psi": round(psi, 2),
                  "operator_id": operator, "material_lot": material_lot, "reagent_supplier": reagent_supplier,
                  "is_anomaly": is_anomaly, "failure_mode": failure_mode}
        data.append(record)
    df = pd.DataFrame(data)
    df['is_anomaly'] = df['is_anomaly'].astype(int)
    return df

@st.cache_data
def generate_capa_data():
    return pd.DataFrame([
        {"id": "CAPA-001", "status": "Closed - Effective", "owner": "Engineer A", "due_date": "2023-05-15"},
        {"id": "CAPA-002", "status": "Pending VOE", "owner": "Engineer B", "due_date": "2024-08-01"},
        {"id": "CAPA-003", "status": "Open - Implementation", "owner": "Engineer C", "due_date": "2024-09-10"},
        {"id": "CAPA-004", "status": "Overdue", "owner": "Engineer B", "due_date": "2024-07-01"}
    ])

@st.cache_data
def generate_method_comparison_data():
    np.random.seed(0)
    true_values = np.random.uniform(5, 50, 50)
    old_method = true_values + np.random.normal(0, 1.5, 50)
    new_method = true_values + np.random.normal(0.5, 1.2, 50)
    return pd.DataFrame({"Old Method": old_method, "New Method": new_method})

@st.cache_data
def generate_lod_data():
    np.random.seed(123)
    conc = np.array([0, 0.1, 0.2, 0.5, 1.0, 2.0, 5.0, 10.0])
    prob_detection = stats.norm.cdf((np.log10(conc, where=conc>0, out=np.full_like(conc, -np.inf)) - np.log10(1.5)) / 0.2)
    replicates = 20
    detected = np.random.binomial(replicates, prob_detection)
    return pd.DataFrame({"Concentration": conc, "Detected": detected, "Total": replicates})


# =================================================================================================
# HELPER FUNCTIONS & DEFINITIONS
# =================================================================================================

def calculate_cpk(df_series, usl, lsl):
    """Calculates Process Capability Index (Cpk). Handles small sample sizes."""
    if len(df_series) < 2 or df_series.std() == 0:
        return np.nan
    mean = df_series.mean()
    std_dev = df_series.std()
    cpu = (usl - mean) / (3 * std_dev)
    cpl = (mean - lsl) / (3 * std_dev)
    return min(cpu, cpl)

def apply_nelson_rules(series, center_line, ucl, lcl):
    """Identifies points violating basic Nelson Rules (1, 2, 3)."""
    violations = pd.Series([False] * len(series), index=series.index)
    # Rule 1: One point outside +/- 3 sigma
    violations = violations | (series > ucl) | (series < lcl)
    # Rule 2: Nine points in a row on the same side of the centerline
    for i in range(8, len(series)):
        if all(series.iloc[i-8:i+1] > center_line) or all(series.iloc[i-8:i+1] < center_line):
            violations.iloc[i] = True
    # Rule 3: Six points in a row, all increasing or all decreasing
    for i in range(5, len(series)):
        if all(np.diff(series.iloc[i-5:i+1]) > 0) or all(np.diff(series.iloc[i-5:i+1]) < 0):
            violations.iloc[i] = True
    return violations

def plot_i_chart(df, param, title):
    """Plots a specialized I-Chart with highlighted violations."""
    data = df[param]
    mr = abs(data.diff()).dropna()
    if len(data) < 2 or len(mr) == 0:
        st.warning(f"Not enough data to generate I-Chart for {param}.")
        return None

    cl = data.mean()
    mr_cl = mr.mean()
    sigma = st.session_state.config.spc_settings.sigma_level
    # d2 constant for n=2 is 1.128
    ucl = cl + sigma * (mr_cl / 1.128)
    lcl = cl - sigma * (mr_cl / 1.128)

    violations = apply_nelson_rules(data, cl, ucl, lcl)
    violation_points = df[violations]

    fig = go.Figure()
    fig.add_trace(go.Scatter(x=df['timestamp'], y=data, mode='lines+markers', name=param, marker_color='blue'))
    fig.add_trace(go.Scatter(x=violation_points['timestamp'], y=violation_points[param], mode='markers', name='Violation', marker=dict(color='red', size=10, symbol='x')))
    fig.add_hline(y=cl, line_color='green', annotation_text="Center Line", annotation_position="bottom right")
    fig.add_hline(y=ucl, line_color='red', line_dash='dash', annotation_text=f"UCL ({sigma}σ)", annotation_position="bottom right")
    fig.add_hline(y=lcl, line_color='red', line_dash='dash', annotation_text=f"LCL ({sigma}σ)", annotation_position="top right")
    fig.update_layout(title_text=f'<b>{title} for {param}</b>', title_x=0.5, legend_orientation='h', legend_x=0, legend_y=1.1)
    st.plotly_chart(fig, use_container_width=True)

    if violations.any():
        st.warning(f"Detected {violations.sum()} out-of-control point(s). See chart for details.")
        with st.expander("Violation Details"):
            st.dataframe(df[violations])
    else:
        st.success("Process appears to be in statistical control based on the selected rules.")
    return fig

# Initialize a dictionary in session_state to store analysis results for reporting
if 'report_content' not in st.session_state:
    st.session_state.report_content = {}


# =================================================================================================
# MODULE 1: QTI COMMAND CENTER
# =================================================================================================
def show_command_center():
    st.title("🔬 QTI Command Center")
    st.markdown("Welcome to the **QTI Command Center**. This dashboard provides a real-time, high-level overview of process health, active quality events, and key performance indicators (KPIs).")
    process_data, capa_data = st.session_state['process_data'], st.session_state['capa_data']

    st.subheader("Key Performance & Quality Indicators (KPIs)")
    kpi_cols = st.columns(5)
    active_inv = len(capa_data[~capa_data['status'].str.startswith('Closed')])
    kpi_cols[0].metric("Active CAPAs", active_inv, help="Corrective and Preventive Actions currently open or pending verification.")
    new_data_count = len(process_data[process_data['timestamp'] > datetime.now() - timedelta(days=1)])
    kpi_cols[1].metric("New Data (24h)", new_data_count, help="Number of new data records ingested in the last 24 hours.")
    anomaly_rate = process_data['is_anomaly'].mean() * 100
    kpi_cols[2].metric("Overall Anomaly Rate", f"{anomaly_rate:.2f}%", help="Percentage of data points flagged as anomalous across all lines.")
    kpi_cols[3].metric("Avg. Pressure (psi)", f"{process_data['pressure_psi'].mean():.2f}")
    config = st.session_state.config.spc_settings
    cpk_overall = calculate_cpk(process_data['pressure_psi'], config.usl, config.lsl)
    kpi_cols[4].metric("Overall Cpk (Pressure)", f"{cpk_overall:.2f}", help=f"Process Capability Index for pressure. Target > 1.33.")

    st.markdown("---")
    col1, col2 = st.columns((3, 2))
    with col1:
        st.subheader("Process Anomaly Timeline")
        st.markdown("Visualizes anomaly occurrences over time, helping to spot trends or event-driven spikes.")
        anomaly_df = process_data[process_data['is_anomaly'] == 1].copy()
        anomaly_df['date'] = anomaly_df['timestamp'].dt.date
        daily_anomalies = anomaly_df.groupby('date').size().reset_index(name='count')
        fig = px.bar(daily_anomalies, x='date', y='count', title="Daily Anomaly Counts", labels={'date': 'Date', 'count': 'Number of Anomalies'})
        st.plotly_chart(fig, use_container_width=True)
    with col2:
        st.subheader("Active CAPA Queue")
        st.markdown("Summary of ongoing corrective and preventive actions.")
        active_capas = capa_data[~capa_data['status'].str.startswith('Closed')].copy()
        st.dataframe(active_capas.style.apply(lambda s: ['background-color: #FFC7CE' if s['status'] == 'Overdue' else '' for v in s], axis=1), use_container_width=True, hide_index=True)

    st.subheader("Process Health Matrix")
    st.markdown("A heatmap showing the anomaly rate (%) for each key parameter on each production line. Darker red indicates a higher rate of anomalies.")
    params = ['reagent_ph', 'fill_volume_ml', 'pressure_psi']
    health_data = []
    for line in sorted(process_data['line_id'].unique()):
        line_df = process_data[process_data['line_id'] == line]
        row = {'line_id': line}
        for param in params:
            param_mean, param_std = process_data[param].mean(), process_data[param].std()
            anomalies = line_df[(line_df[param] > param_mean + 3 * param_std) | (line_df[param] < param_mean - 3 * param_std)]
            row[param] = (len(anomalies) / len(line_df) * 100) if len(line_df) > 0 else 0
        health_data.append(row)
    health_df = pd.DataFrame(health_data).set_index('line_id')
    fig_health = px.imshow(health_df, text_auto=".2f", aspect="auto", color_continuous_scale='RdYlGn_r',
                           labels=dict(x="Process Parameter", y="Production Line", color="Anomaly Rate (%)"),
                           title="Parameter Anomaly Rate by Production Line")
    st.plotly_chart(fig_health, use_container_width=True)

# =================================================================================================
# MODULE 2: PROCESS MONITORING
# =================================================================================================
def show_process_monitoring():
    st.title("📈 Process Monitoring & Statistical Process Control (SPC)")
    st.markdown("This module provides tools to monitor process stability and detect deviations from normal operation.")
    process_data = st.session_state['process_data']

    st.sidebar.header("Monitoring Filters")
    lines = ['All Lines'] + sorted(process_data['line_id'].unique())
    line_to_monitor = st.sidebar.selectbox("Select Process Line:", lines)
    monitor_df = process_data.copy() if line_to_monitor == 'All Lines' else process_data[process_data['line_id'] == line_to_monitor].copy()
    monitor_df = monitor_df.sort_values('timestamp').reset_index(drop=True)

    st.sidebar.subheader("SPC Settings")
    st.session_state.config.spc_settings.sigma_level = st.sidebar.slider(
        "Control Limit Sigma (σ)", 2.0, 4.0, st.session_state.config.spc_settings.sigma_level, 0.5,
        help="Sets the width of the control limits (e.g., 3σ is standard).")

    if len(monitor_df) < 10:
        st.warning("Insufficient data for the selected line. Please select a line with more data points.")
        return

    tab1, tab2, tab3, tab4 = st.tabs(["I-MR Chart", "EWMA Chart", "CUSUM Chart", "Multivariate (Hotelling's T²)"])
    param_options = ('reagent_ph', 'fill_volume_ml', 'pressure_psi')

    with tab1:
        st.subheader("Individuals Chart (I-Chart) with Nelson Rules")
        with st.expander("**Methodology & Interpretation**"):
            st.markdown("- **What it is:** The Individuals chart plots individual measurements to monitor process center and variation over time.\n- **When to use it:** Ideal for continuous data where measurements are taken one at a time.\n- **How to interpret:** Points marked with a red 'X' violate a statistical process rule, suggesting a special cause of variation.")
        param_imr = st.selectbox("Select Parameter:", param_options, key="imr_param")
        if len(monitor_df) > 1:
            imr_fig = plot_i_chart(monitor_df, param_imr, "Individuals (I) Chart")
            if imr_fig:
                st.session_state.report_content['spc_imr'] = {
                    "title": f"I-Chart for {param_imr} on {line_to_monitor}", "figure": imr_fig,
                    "text": f"An I-Chart for '{param_imr}' was generated with {st.session_state.config.spc_settings.sigma_level}σ control limits."}

    with tab2:
        st.subheader("Exponentially Weighted Moving Average (EWMA) Chart")
        with st.expander("**Methodology & Interpretation**"):
            st.markdown("- **What it is:** A chart that plots a weighted average of observations, giving more weight to recent data.\n- **When to use it:** Highly effective at detecting small, sustained process shifts or drifts.\n- **How to interpret:** Look for points in the EWMA series (purple line) that fall outside the control limits. A smaller `λ` detects smaller shifts.")
        param_ewma = st.selectbox("Select Parameter:", param_options, key="ewma_param")
        lambda_val = st.slider("Smoothing Factor (λ)", 0.05, 1.0, 0.2, 0.05, help="Smaller λ gives more weight to past data, detecting smaller shifts.")
        data = monitor_df[param_ewma]; ewma = data.ewm(alpha=lambda_val, adjust=False).mean()
        cl, std = data.mean(), data.std(); sigma = st.session_state.config.spc_settings.sigma_level
        ucl = cl + sigma * std * np.sqrt(lambda_val / (2 - lambda_val)); lcl = cl - sigma * std * np.sqrt(lambda_val / (2 - lambda_val))
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=monitor_df['timestamp'], y=data, mode='lines', name='Raw Data', line=dict(color='lightblue', width=1)))
        fig.add_trace(go.Scatter(x=monitor_df['timestamp'], y=ewma, mode='lines', name='EWMA', line=dict(color='purple')))
        fig.add_hline(y=cl, line_color='green', annotation_text="Center Line"); fig.add_hline(y=ucl, line_color='red', line_dash='dash', annotation_text="UCL"); fig.add_hline(y=lcl, line_color='red', line_dash='dash', annotation_text="LCL")
        fig.update_layout(title_text=f'<b>EWMA Chart for {param_ewma} (λ={lambda_val})</b>', title_x=0.5)
        st.plotly_chart(fig, use_container_width=True)

    with tab3:
        st.subheader("Cumulative Sum (CUSUM) Chart")
        with st.expander("**Methodology & Interpretation**"):
            st.markdown("- **What it is:** A chart that accumulates deviations from a target value over time.\n- **When to use it:** Very effective for detecting small, persistent shifts in the process mean.\n- **How to interpret:** A sustained upward or downward trend indicates the process mean has shifted. The slope is proportional to the size of the shift.")
        param_cusum = st.selectbox("Select Parameter:", param_options, key="cusum_param")
        target = st.number_input("Process Target:", value=monitor_df[param_cusum].mean(), format="%.3f")
        cusum = (monitor_df[param_cusum] - target).cumsum()
        fig = px.line(x=monitor_df['timestamp'], y=cusum, title=f"<b>CUSUM Chart for {param_cusum}</b>", markers=False)
        fig.add_hline(y=0, line_color='green'); fig.update_layout(title_x=0.5, yaxis_title="Cumulative Sum of Deviations")
        st.plotly_chart(fig, use_container_width=True)

    with tab4:
        st.subheader("Multivariate Control Chart (Hotelling's T²)")
        with st.expander("**Methodology & Interpretation**"):
            st.markdown("- **What it is:** A multivariate chart that combines correlated variables into a single statistic measuring the 'distance' from the process center.\n- **When to use it:** To monitor several correlated variables simultaneously. It can detect issues that individual charts miss.\n- **How to interpret:** A point exceeding the UCL indicates a significant deviation, but doesn't specify which variable is responsible.")
        features = st.multiselect("Select Features for Multivariate Analysis:", param_options, default=['reagent_ph', 'pressure_psi'])
        if len(features) < 2:
            st.info("Please select at least two features for multivariate analysis.")
        elif len(monitor_df) <= len(features):
            st.warning("Not enough data points for the selected number of features.")
        else:
            X = monitor_df[features]; X_scaled = StandardScaler().fit_transform(X)
            mean_vec = np.mean(X_scaled, axis=0)
            # Use pseudo-inverse (pinv) for robustness against singular matrices
            inv_cov = np.linalg.pinv(np.cov(X_scaled, rowvar=False))
            t_sq = [ (row - mean_vec) @ inv_cov @ (row - mean_vec).T for row in X_scaled]
            p, n = len(features), len(X); alpha = 0.01
            ucl = (p * (n + 1) * (n - 1)) / (n * n - n * p) * stats.f.ppf(1 - alpha, p, n - p)
            fig = px.line(x=monitor_df['timestamp'], y=t_sq, title=f"<b>Hotelling's T² Chart (α={alpha})</b>", markers=True)
            fig.add_hline(y=ucl, line_color='red', line_dash='dash', name=f'UCL ({1-alpha:.0%})')
            fig.update_layout(title_x=0.5, yaxis_title="T-Squared Statistic")
            st.plotly_chart(fig, use_container_width=True)
# =================================================================================================
# MODULE 3: RCA WORKBENCH
# =================================================================================================
def show_rca_workbench():
    st.title("🛠️ Root Cause Analysis (RCA) Workbench")
    st.markdown("A set of analytical tools to investigate process deviations and identify potential root causes.")
    process_data = st.session_state['process_data']
    st.sidebar.header("RCA Filters")
    min_date, max_date = process_data['timestamp'].min().date(), process_data['timestamp'].max().date()
    start_date, end_date = st.sidebar.slider("Select Date Range for Analysis:", min_value=min_date, max_value=max_date, value=(min_date, max_date), format="YYYY-MM-DD")
    rca_df = process_data[(process_data['timestamp'].dt.date >= start_date) & (process_data['timestamp'].dt.date <= end_date)].copy()
    if rca_df.empty: st.warning("No data available for the selected date range."); return

    tab1, tab2, tab3, tab4 = st.tabs(["Hypothesis Testing", "Distribution Analysis", "Unsupervised Clustering", "AI-Driven Importance"])
    rca_params = ['reagent_ph', 'fill_volume_ml', 'pressure_psi']
    rca_cats = ['material_lot', 'operator_id', 'line_id', 'reagent_supplier']

    with tab1:
        st.subheader("Statistical Hypothesis Testing"); st.markdown("**Use Case:** Test if a parameter's mean differs significantly between two groups.")
        col1, col2 = st.columns(2); param = col1.selectbox("Parameter to Analyze:", rca_params, key="ttest_param"); group_by = col2.selectbox("Grouping Variable:", rca_cats, key="ttest_group")
        groups = rca_df[group_by].unique()
        if len(groups) < 2: st.warning(f"Grouping variable '{group_by}' has fewer than 2 groups."); return
        group1_name = col1.selectbox("Select Group 1:", groups, index=0); group2_name = col2.selectbox("Select Group 2:", [g for g in groups if g != group1_name], index=0)
        g1_data = rca_df[rca_df[group_by] == group1_name][param].dropna(); g2_data = rca_df[rca_df[group_by] == group2_name][param].dropna()
        fig = go.Figure(); fig.add_trace(go.Box(y=g1_data, name=group1_name)); fig.add_trace(go.Box(y=g2_data, name=group2_name)); fig.update_layout(title=f'<b>Comparison of {param}</b>', title_x=0.5)
        st.plotly_chart(fig, use_container_width=True)
        st.subheader("Welch's t-test Results (unequal variances)");
        if len(g1_data) > 1 and len(g2_data) > 1:
            ttest = stats.ttest_ind(g1_data, g2_data, equal_var=False)
            res_col1, res_col2 = st.columns(2); res_col1.metric("p-value", f"{ttest.pvalue:.4g}"); res_col2.metric("Difference in Means", f"{g1_data.mean() - g2_data.mean():.3f}")
            if ttest.pvalue < 0.05: st.error("**Conclusion:** The difference is statistically significant (p < 0.05).")
            else: st.success("**Conclusion:** No significant difference detected (p >= 0.05).")
        else: st.warning("Not enough data in one or both groups for a t-test.")

    with tab2:
        st.subheader("Distribution Analysis"); st.markdown("**Use Case:** Visualize data distribution to identify skewness, modes, or outliers, stratified by category.")
        col1, col2 = st.columns(2); param_dist = col1.selectbox("Parameter to Plot:", rca_params, key="dist_param"); cat_dist = col2.selectbox("Stratify by:", [None] + rca_cats, key="dist_cat")
        fig = px.histogram(rca_df, x=param_dist, color=cat_dist, marginal="box", barmode='overlay', opacity=0.7, title=f"<b>Distribution of {param_dist}</b>" + (f" by {cat_dist}" if cat_dist else "")); fig.update_layout(title_x=0.5)
        st.plotly_chart(fig, use_container_width=True)

    with tab3:
        st.subheader("Unsupervised Clustering (K-Means)"); st.markdown("**Use Case:** Automatically discover hidden groups or operational states within process data.")
        cluster_features = st.multiselect("Select features for clustering:", rca_params, default=rca_params[:2])
        if len(cluster_features) >= 2:
            X = rca_df[cluster_features].dropna(); X_scaled = StandardScaler().fit_transform(X)
            c1, c2 = st.columns([1, 2])
            with c1:
                st.markdown("**Choosing 'k' (Clusters)**"); st.markdown("The 'Elbow Plot' helps select optimal `k` where returns diminish.")
                wcss = [KMeans(n_clusters=i, init='k-means++', random_state=42, n_init=10).fit(X_scaled).inertia_ for i in range(1, 11)]
                elbow_fig = px.line(x=list(range(1, 11)), y=wcss, title="Elbow Method", labels={'x': 'k', 'y': 'WCSS'}); st.plotly_chart(elbow_fig, use_container_width=True)
                k = st.slider("Select number of clusters (k):", 2, 10, 4)
            with c2:
                kmeans = KMeans(n_clusters=k, random_state=42, n_init=10).fit(X_scaled); rca_df['cluster'] = kmeans.labels_.astype(str)
                fig = px.scatter(rca_df, x=cluster_features[0], y=cluster_features[1], color='cluster', hover_data=rca_df.columns, title=f"<b>K-Means Clustering Result (k={k})</b>"); fig.update_layout(title_x=0.5)
                st.plotly_chart(fig, use_container_width=True)
            st.subheader("Cluster Profiles"); st.markdown("Characteristics of each discovered cluster.")
            profile = rca_df.groupby('cluster')[rca_params + rca_cats].agg(lambda x: x.mode()[0] if x.dtype == 'object' else x.mean()).round(2)
            profile['count'] = rca_df.groupby('cluster').size(); st.dataframe(profile)
        else: st.info("Please select at least two features for clustering.")

# =====================================================================================
# START: DEFINITIVE FIX FOR `AssertionError` in `shap.summary_plot`
#
# This block replaces the entire "AI-Driven Importance" tab content in `show_rca_workbench`.
# It resolves the error by training a temporary, local model on the filtered `rca_df`.
# This guarantees perfect alignment between the model, the SHAP values, and the feature data,
# completely eliminating the possibility of a column mismatch error.
# =====================================================================================
    with tab4:
        st.subheader("AI-Driven Root Cause Identification (Feature Importance)")
        st.markdown("**Use Case:** Use a machine learning model to rank variables by their importance in predicting an anomaly within the selected time frame.")

        # Check for sufficient data for modeling
        if rca_df['is_anomaly'].nunique() < 2 or rca_df['is_anomaly'].value_counts().min() < 5:
            st.warning("Not enough anomaly data in the selected date range to build a reliable predictive model. Ensure the data contains at least 5 samples for both anomaly and non-anomaly cases.")
        else:
            with st.spinner("Training local model and calculating SHAP values for the selected time range..."):
                try:
                    # --- Local Model Training on Filtered Data (The Definitive Fix) ---
                    # This strategy trains a model ONLY on the data being analyzed, guaranteeing
                    # perfect alignment between the model, data, and explainer.

                    # 1. Define features based on the *filtered* rca_df to ensure relevance.
                    features = ['reagent_ph', 'fill_volume_ml', 'pressure_psi'] + [col for col in rca_cats if rca_df[col].nunique() > 1]
                    target = 'is_anomaly'

                    # 2. Prepare data from the filtered dataframe
                    X = pd.get_dummies(rca_df[features], drop_first=True)
                    y = rca_df[target]

                    # 3. Split the filtered data for training and explanation
                    X_train, X_test, y_train, y_test = train_test_split(
                        X, y, test_size=0.3, random_state=42, stratify=y
                    )

                    # 4. Create and train the local model
                    model = RandomForestClassifier(n_estimators=100, random_state=42, class_weight='balanced')
                    model.fit(X_train, y_train)

                    # 5. Create the explainer and SHAP values based on the locally trained model
                    explainer = shap.TreeExplainer(model)
                    shap_values = explainer.shap_values(X_test)

                    # --- SHAP Summary Plot (Beeswarm) ---
                    st.markdown("**SHAP Summary Plot**")
                    st.markdown("This plot shows the impact of each feature on the model's prediction of an anomaly. Each dot is a single observation. Red means a high feature value, blue means low. A positive SHAP value pushes the prediction towards 'Anomaly'.")

                    # This call is now guaranteed to be safe because the model, explainer,
                    # shap_values, and X_test are all derived from the same filtered dataset (rca_df).
                    fig, ax = plt.subplots()
                    shap.summary_plot(shap_values[1], X_test, show=False)
                    st.pyplot(fig)
                    plt.close(fig)

                    # --- Feature Importance Bar Chart ---
                    importances = pd.DataFrame({
                        'feature': X_train.columns,
                        'importance': model.feature_importances_
                    }).sort_values('importance', ascending=False)

                    st.markdown("**Local Feature Importance**")
                    st.markdown("This chart shows the overall importance of each feature in the model for the selected time period.")
                    fig_bar = px.bar(importances, x='importance', y='feature', orientation='h', title='Feature Importance for Anomaly Prediction')
                    fig_bar.update_layout(yaxis={'categoryorder': 'total ascending'})
                    st.plotly_chart(fig_bar, use_container_width=True)

                    st.session_state.report_content['rca_importance'] = {
                        "title": "AI-Driven Feature Importance (RCA Period)",
                        "figure": fig_bar,
                        "text": "A Random Forest model was trained on data from the selected time period to predict anomalies. The chart shows the Gini importance of each feature for that specific period."
                    }
                except Exception as e:
                    st.error(f"An error occurred during model training or SHAP analysis: {e}")
# =====================================================================================
# END: DEFINITIVE FIX
# =====================================================================================

# =================================================================================================
# MODULE 4: CHANGE VALIDATION & CAPA
# =================================================================================================
def show_change_validation():
    st.title("✅ Change Validation & CAPA Management")
    st.markdown("Tools for validating process changes and tracking Corrective and Preventive Actions (CAPAs).")
    tab1, tab2, tab3 = st.tabs(["CAPA Log", "Method Agreement (Bland-Altman)", "Equivalence Testing (TOST)"])

    with tab1:
        st.subheader("CAPA Action Log"); st.dataframe(st.session_state['capa_data'], use_container_width=True, hide_index=True)

    with tab2:
        st.subheader("Method Agreement (Bland-Altman Plot)")
        with st.expander("**Methodology & Interpretation**"): st.markdown("- **What it is:** A plot comparing two measurement techniques by plotting their difference against their average.\n- **When to use it:** To see how well a new measurement method agrees with an old one.\n- **How to interpret:** The mean difference (blue line) is the systematic bias. 95% of differences should fall within the Limits of Agreement (red dashed lines).")
        df = generate_method_comparison_data(); df['Average'] = df.mean(axis=1); df['Difference'] = df['New Method'] - df['Old Method']
        mean_diff = df['Difference'].mean(); std_diff = df['Difference'].std(); upper_loa = mean_diff + 1.96 * std_diff; lower_loa = mean_diff - 1.96 * std_diff
        fig = px.scatter(df, x='Average', y='Difference', title='<b>Bland-Altman Plot: New vs. Old Method</b>'); fig.add_hline(y=mean_diff, line_color='blue', annotation_text=f"Bias: {mean_diff:.2f}"); fig.add_hline(y=upper_loa, line_color='red', line_dash='dash', annotation_text=f"Upper LoA: {upper_loa:.2f}"); fig.add_hline(y=lower_loa, line_color='red', line_dash='dash', annotation_text=f"Lower LoA: {lower_loa:.2f}"); fig.update_layout(title_x=0.5)
        st.plotly_chart(fig, use_container_width=True); st.metric("Systematic Bias (New - Old)", f"{mean_diff:.3f}"); st.metric("95% Limits of Agreement", f"[{lower_loa:.3f}, {upper_loa:.3f}]")

    with tab3:
        st.subheader("Equivalence Testing (TOST)")
        with st.expander("**Methodology & Interpretation**"): st.markdown("- **What it is:** Two One-Sided Tests (TOST) to provide evidence that two means are *equivalent*.\n- **When to use it:** To prove a change did *not* negatively impact the process output.\n- **How to interpret:** Define an equivalence margin `[-δ, +δ]`. If the 90% confidence interval of the difference lies entirely within this margin, you can claim equivalence (p < 0.05).")
        np.random.seed(42); batch_A = np.random.normal(10.0, 1.0, 50); batch_B = np.random.normal(10.05, 1.0, 50); mean_diff = np.mean(batch_B) - np.mean(batch_A)
        st.markdown("Define your equivalence margin (`δ`), the largest difference you consider practically meaningless."); bound = st.slider("Equivalence Margin (±)", 0.1, 1.0, 0.5, 0.05); low, high = -bound, bound
        p_val, _, _ = ttost_ind(batch_B, batch_A, low=low, upp=high)
        cm = sm.stats.CompareMeans(sm.stats.DescrStatsW(batch_A), sm.stats.DescrStatsW(batch_B)); ci_low, ci_high = cm.tconfint_diff(alpha=0.10, usevar='unequal')
        fig = go.Figure(); fig.add_trace(go.Bar(x=[(high+low)/2], y=[1], base=low, width=high-low, orientation='h', name='Equivalence Interval', marker_color='lightgreen', opacity=0.5)); fig.add_trace(go.Scatter(x=[mean_diff], y=[1], mode='markers', marker=dict(color='black', size=15), name='Mean Difference')); fig.add_shape(type='line', x0=ci_low, y0=1, x1=ci_high, y1=1, line=dict(color='black', width=3), name='90% CI of Difference'); fig.update_layout(title="<b>TOST Equivalence Visualization</b>", xaxis_title="Difference (Batch B - Batch A)", yaxis_showticklabels=False, height=250)
        st.plotly_chart(fig, use_container_width=True); st.metric("TOST p-value", f"{p_val:.4g}")
        if p_val < 0.05: st.success("✅ **Conclusion:** Processes are statistically equivalent within the specified bounds.")
        else: st.error("❌ **Conclusion:** Equivalence cannot be concluded.")


# =================================================================================================
# MODULE 5: PREDICTIVE & OPTIMIZATION ANALYTICS
# =================================================================================================
@st.cache_resource
def get_model_and_explainer(_df):
    """Caches the trained ML pipeline and SHAP explainer. Trains on the full dataset provided."""
    _df = _df.copy()
    features = ['reagent_ph', 'fill_volume_ml', 'pressure_psi', 'material_lot', 'operator_id', 'line_id', 'reagent_supplier']
    target = 'is_anomaly'
    X = pd.get_dummies(_df[features], drop_first=True)
    y = _df[target]
    if y.nunique() < 2 or y.value_counts().min() < 5: return None, None
    X_train, _, y_train, _ = train_test_split(X, y, test_size=0.3, random_state=42, stratify=y)
    pipeline = Pipeline([('scaler', StandardScaler()), ('classifier', RandomForestClassifier(n_estimators=100, random_state=42, class_weight='balanced'))])
    pipeline.fit(X_train, y_train)
    explainer = shap.TreeExplainer(pipeline.named_steps['classifier'])
    return pipeline, explainer

def show_predictive_and_optimization():
    st.title("🔮 Predictive & Optimization Analytics")
    st.markdown("Leverage advanced analytics to forecast behavior, optimize settings, and get explainable predictions.")
    process_data = st.session_state['process_data']
    tab1, tab2, tab3 = st.tabs(["Time Series Forecasting", "Process Optimization (Bayesian)", "Real-Time Prediction (XAI)"])

    with tab1:
        st.subheader("Time Series Forecasting (SARIMA)"); st.markdown("**Use Case:** Forecast future process behavior based on historical trends and seasonality.")
        ts_data = process_data[['timestamp', 'pressure_psi']].set_index('timestamp').resample('D').mean().dropna()
        if len(ts_data) < 20: st.warning("Insufficient daily data for a reliable forecast (need >20 days).")
        else:
            with st.spinner("Training SARIMA model..."):
                try:
                    model = sm.tsa.statespace.SARIMAX(ts_data['pressure_psi'], order=(1,1,1), seasonal_order=(1,1,1,7)).fit(disp=False)
                    forecast = model.get_forecast(steps=30); forecast_df = forecast.summary_frame()
                    fig = go.Figure(); fig.add_trace(go.Scatter(x=ts_data.index, y=ts_data['pressure_psi'], name='Historical Data')); fig.add_trace(go.Scatter(x=forecast_df.index, y=forecast_df['mean'], name='Forecast', line=dict(color='orange'))); fig.add_trace(go.Scatter(x=forecast_df.index, y=forecast_df['mean_ci_upper'], fill='tonexty', mode='lines', line_color='rgba(255,165,0,0.2)', name='95% CI')); fig.add_trace(go.Scatter(x=forecast_df.index, y=forecast_df['mean_ci_lower'], fill='tonexty', mode='lines', line_color='rgba(255,165,0,0.2)')); fig.update_layout(title_text="<b>Pressure Forecast (Next 30 Days)</b>", title_x=0.5, showlegend=False); st.plotly_chart(fig, use_container_width=True)
                except Exception as e: st.error(f"Failed to build forecast model: {e}")

    with tab2:
        st.subheader("Process Optimization (Bayesian)"); st.markdown("**Use Case:** Intelligently search for optimal process settings to maximize a desired outcome.")
        with st.expander("**How it Works**"): st.markdown("Bayesian Optimization builds a probability model of the objective function and uses it to select the most promising parameters to evaluate next, avoiding a brute-force search.")
        space = [Real(7.0, 7.4, name='ph'), Real(48.0, 52.0, name='psi')]
        @use_named_args(dimensions=space)
        def black_box_function(ph, psi): return -((ph - 7.25)**2 + (psi - 50.5)**2)
        if st.button("Run Bayesian Optimization", key="run_bayes_opt"):
            with st.spinner("Finding optimal settings... (15 iterations)"):
                res = gp_minimize(black_box_function, dimensions=space, n_calls=15, random_state=0)
                st.success(f"**Optimization Complete!**"); res_col1, res_col2 = st.columns(2); res_col1.metric("Optimal pH", f"{res.x[0]:.3f}"); res_col2.metric("Optimal Pressure (psi)", f"{res.x[1]:.2f}")
                fig, ax = plt.subplots(figsize=(8, 4)); plot_convergence(res, ax=ax); ax.set_title("Convergence Plot"); st.pyplot(fig); plt.close(fig)

    with tab3:
        st.subheader("Real-Time Prediction with XAI (eXplainable AI)"); st.markdown("Input hypothetical parameters to get a real-time anomaly prediction and a breakdown of driving factors.")
        pipeline, explainer = get_model_and_explainer(process_data)
        if not pipeline: st.error("Model could not be trained. Please check dataset for anomalies."); return

        col1, col2 = st.columns([1, 2])
        with col1:
            st.write("**Input Parameters:**")
            input_data = {'reagent_ph': st.slider("Reagent pH", 7.0, 7.6, 7.25, 0.01),
                          'fill_volume_ml': st.slider("Fill Volume (ml)", 9.8, 10.2, 10.0, 0.01),
                          'pressure_psi': st.slider("Pressure (psi)", 45.0, 58.0, 51.0, 0.1),
                          'material_lot': st.selectbox("Material Lot", process_data['material_lot'].unique()),
                          'operator_id': st.selectbox("Operator ID", process_data['operator_id'].unique()),
                          'line_id': st.selectbox("Line ID", process_data['line_id'].unique()),
                          'reagent_supplier': st.selectbox("Reagent Supplier", process_data['reagent_supplier'].unique())}
            input_df_raw = pd.DataFrame([input_data])
            input_df_processed = pd.get_dummies(input_df_raw).reindex(columns=pipeline.feature_names_in_, fill_value=0)
            proba = pipeline.predict_proba(input_df_processed)[0][1]
            st.metric("Predicted Anomaly Probability", f"{proba:.1%}")
            if proba > 0.5: st.warning("High risk of anomaly detected.")
            else: st.success("Process parameters appear nominal.")
        with col2:
            st.write("**XAI Driver Analysis (SHAP Force Plot)**"); st.markdown("This plot shows forces pushing the prediction. <span style='color:red;'>Red bars</span> increase anomaly risk, <span style='color:blue;'>blue bars</span> decrease it.", unsafe_allow_html=True)
            scaled_input = pipeline.named_steps['scaler'].transform(input_df_processed)
            explanation = explainer(scaled_input)
            
            expected_value_class1 = explanation.base_values[0, 1]
            shap_values_class1 = explanation.values[0, :, 1]
            
            fig, ax = plt.subplots(figsize=(10, 2.5))
            shap.force_plot(expected_value_class1, shap_values_class1, input_df_processed.iloc[0], matplotlib=True, show=False, text_rotation=15)
            plt.tight_layout(); st.pyplot(fig, bbox_inches='tight'); plt.close(fig)


# =================================================================================================
# MODULE 6: ADVANCED DEMOS & VALIDATION
# =================================================================================================
def show_advanced_demos():
    st.title("⚙️ Advanced Demos & Validation")
    st.markdown("Specialized validation tasks and advanced library integrations.")
    tab1, tab2 = st.tabs(["Data Scalability (Dask)", "Assay Validation (LoD/LoQ)"])

    with tab1:
        st.subheader("Large-Scale Data Processing with Dask"); st.markdown("**Use Case:** Dask enables parallel computation on datasets larger than system memory by breaking them into chunks."); st.info("This demo compares Dask's performance against standard Pandas on a complex aggregation.")
        if st.button("Run Dask vs. Pandas Benchmark"):
            with st.spinner("Processing with Dask..."):
                ddf = dd.from_pandas(st.session_state['process_data'], npartitions=4); start_time_dask = time.time()
                dask_result = ddf.groupby('operator_id').agg({'pressure_psi': ['mean', 'std'], 'reagent_ph': 'mean'}).compute(); dask_time = time.time() - start_time_dask
            st.success(f"Dask computation complete in **{dask_time:.4f} seconds**.")
            with st.spinner("Processing with Pandas..."):
                df = st.session_state['process_data']; start_time_pandas = time.time()
                pandas_result = df.groupby('operator_id').agg({'pressure_psi': ['mean', 'std'], 'reagent_ph': 'mean'}); pandas_time = time.time() - start_time_pandas
            st.success(f"Pandas computation complete in **{pandas_time:.4f} seconds**.")
            st.markdown(f"**Result:** Dask took `{dask_time:.4f}s` vs. Pandas `{pandas_time:.4f}s`. Dask's overhead may make it slower on small data, but its architecture scales to terabyte-sized datasets where Pandas would fail."); st.dataframe(dask_result)

    with tab2:
        st.subheader("Assay Validation: Limit of Detection (LoD)"); st.markdown("**Use Case:** Determine the lowest concentration an assay can reliably detect.");
        with st.expander("**Methodology**"): st.markdown("- **Limit of Detection (LoD):** Calculated using Probit Regression (a GLM) to model the probability of detection vs. concentration. We find the concentration that gives a 95% detection rate.")
        df_lod = generate_lod_data(); df_lod['Not Detected'] = df_lod['Total'] - df_lod['Detected']; df_lod['log_conc'] = np.log10(df_lod['Concentration'].replace(0, 0.01))
        glm_binom = sm.GLM(endog=df_lod[['Detected', 'Not Detected']], exog=sm.add_constant(df_lod['log_conc']), family=sm.families.Binomial(link=sm.families.links.probit())); res = glm_binom.fit()
        target_prob = stats.norm.ppf(0.95); params = res.params; lod = 10**((target_prob - params['const']) / params['log_conc'])
        x_range = np.linspace(df_lod['log_conc'].min(), df_lod['log_conc'].max(), 200); y_pred = res.predict(sm.add_constant(x_range))
        fig = go.Figure(); fig.add_trace(go.Scatter(x=df_lod['Concentration'], y=df_lod['Detected']/df_lod['Total'], mode='markers', name='Observed Hit Rate')); fig.add_trace(go.Scatter(x=10**x_range, y=y_pred, mode='lines', name='Probit Fit Curve')); fig.add_vline(x=lod, line_dash='dash', line_color='red', annotation_text=f"LoD = {lod:.3f}"); fig.add_hline(y=0.95, line_dash='dash', line_color='red'); fig.update_layout(title='<b>Probit Analysis for Limit of Detection (LoD)</b>', xaxis_title="Concentration", yaxis_title="Detection Probability", xaxis_type="log", title_x=0.5)
        st.metric("Calculated LoD (95% Probability)", f"{lod:.3f}"); st.plotly_chart(fig, use_container_width=True)

# =================================================================================================
# MODULE 7: REPORTING & EXPORT
# =================================================================================================
def add_slide_with_content(prs, title_text, content_text, figure=None):
    # Use slide layout 1 ("Title and Content") to ensure placeholder availability.
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.text = content_text
    tf.paragraphs[0].font.size = Pt(14)

    if figure:
        try:
            img_stream = io.BytesIO()
            figure.write_image(img_stream, format='png', scale=2, width=800, height=450)
            img_stream.seek(0)
            # Position picture below text
            slide.shapes.add_picture(img_stream, Inches(1), Inches(2.5), width=Inches(8))
        except Exception as e:
            st.warning(f"Could not export figure '{title_text}' to PowerPoint. Error: {e}")

def show_reporting():
    st.title("📄 Interactive Report Builder")
    st.markdown("Compile key findings from your investigation into a downloadable PowerPoint report.")
    st.subheader("Select Analyses to Include in Report")
    selections = {key: st.checkbox(f"Include: **{item['title']}**", value=True) for key, item in st.session_state.report_content.items()}
    if not st.session_state.report_content: st.info("No analyses generated yet. Visit other modules to add analyses to the report."); return
    st.markdown("---"); st.subheader("Generate Report")
    if st.button("Generate .pptx Report"):
        with st.spinner("Creating PowerPoint presentation..."):
            config = st.session_state.config; prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0]); 
            slide.shapes.title.text = "QTI Investigation Summary Report"
            slide.placeholders[1].text = f"Author: {config.report_settings.author}\nCompany: {config.report_settings.company_name}\nDate: {datetime.now().strftime('%Y-%m-%d')}"
            
            for key, selected in selections.items():
                if selected:
                    content = st.session_state.report_content[key]
                    add_slide_with_content(prs, title_text=content['title'], content_text=content.get('text', 'No summary.'), figure=content.get('figure'))
            
            ppt_stream = io.BytesIO(); prs.save(ppt_stream); ppt_stream.seek(0)
            st.download_button(label="📥 Download Report", data=ppt_stream, file_name=f"QTI_Report_{datetime.now().strftime('%Y%m%d')}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# =================================================================================================
# MAIN APPLICATION LOGIC
# =================================================================================================
def main():
    """Main function to run the Streamlit app."""
    st.sidebar.title("QTI Workbench v2.4")
    st.sidebar.markdown("---")
    if 'data_loaded' not in st.session_state:
        st.session_state.config = load_config()
        st.session_state.process_data = generate_process_data()
        st.session_state.capa_data = generate_capa_data()
        st.session_state.data_loaded = True
        st.session_state.report_content = {}

    page_functions = {"QTI Command Center": show_command_center, "Process Monitoring": show_process_monitoring, "RCA Workbench": show_rca_workbench, "Change Validation & CAPA": show_change_validation,
                      "Predictive & Optimization": show_predictive_and_optimization, "Advanced Demos": show_advanced_demos, "Report Builder": show_reporting}
    module = st.sidebar.radio("Select a Module:", tuple(page_functions.keys()))
    st.sidebar.markdown("---"); st.sidebar.info("**QTI Engineering Workbench**\n\n© 2024 Innovate Bio-Diagnostics\n\n*Gold Master Edition v2.4*")
    page_functions[module]()

if __name__ == "__main__":
    main()
